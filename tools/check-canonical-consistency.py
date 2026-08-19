#!/usr/bin/env python3
"""
Mechanical consistency checker for "Miksi tekoälyyn investoidaan biljoonia?"
(Finnish translation of "Why Are They Spending Trillions on AI?" — adapted
from the English repo's tools/check-canonical-consistency.py. Figure-matching
is adjusted for this repo's Finnish-locale number formatting, comma decimal
separators, per TERMINOLOGY.md.)

Purpose: no chart, slide, PDF, or workbook in this package should be able to
quietly diverge from data/canonical-cost-model.csv, which is this package's
single source of truth for tier cost figures. This script:

  1. Loads data/canonical-cost-model.csv and DERIVES its expected-figure
     strings from it at runtime (it does not hard-code the canonical numbers
     anywhere else in this file) — editing the CSV changes what the checker
     looks for, without needing to edit this script.
  2. Confirms every canonical low/high figure appears verbatim in the
     whitepaper and in every other markdown asset that is expected to cite
     tier costs.
  3. Confirms a fixed list of previously-found-and-fixed STALE figures
     (superseded ranges that do not match the canonical table) does not
     reappear anywhere outside an explicit "Corrected"/"Resolved"
     self-correction note — this package's established pattern for
     discussing a superseded number by name.
  4. Extracts text from the built PDF (01-whitepaper.pdf) and checks it the
     same way as the whitepaper markdown, so a stale PDF that wasn't
     regenerated after a markdown fix is caught, not just a stale markdown.
  5. Opens the companion Excel workbook and checks specific regression-prone
     cells (a hyperscale utilization assumption, two formula-string cells).
  6. Opens the slide deck and checks its rendered text the same way.

FAIL CLOSED: if a required dependency (openpyxl, python-pptx, pypdf) is not
importable, or a required file is missing, that check is reported as a
FAILURE, not silently skipped. This script never prints PASS while having
skipped part of its own stated scope.

Setup (from repo root):
    python3 -m venv .venv && source .venv/bin/activate
    pip install -r tools/requirements.txt
    python3 tools/check-canonical-consistency.py

Run: python3 tools/check-canonical-consistency.py
Exit code 0 = clean. Exit code 1 = drift found, or a required dependency/file
is missing (see printed report either way).
"""
import csv
import importlib
import os
import re
import sys

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CANONICAL_CSV = os.path.join(ROOT, "data", "canonical-cost-model.csv")

REQUIRED_DEPS = ["openpyxl", "pptx", "pypdf"]

# Assets that are expected to cite tier cost figures and are checked against
# the canonical CSV plus the stale-figure blocklist below.
MARKDOWN_FILES_TO_CHECK = [
    "01-whitepaper.md",
    "05-workbook-token-factory-scenarios.md",
    "13-slide-deck-outline.md",
    "17-visual-asset-briefs.md",
]

# Regex fragments for figures that were found stale and fixed in this
# package's history. Each entry: (regex, human-readable description). A line
# containing one of these is a FAIL unless it also contains a self-correction
# marker (SELF_CORRECTION_MARKERS below) on the same line.
# Finnish-locale variants (comma decimal separator) of the same stale
# figures, since this repo's prose uses natural Finnish number formatting
# (e.g. "1,37" not "1.37") per TERMINOLOGY.md's Number formatting section.
STALE_PATTERNS = [
    (r"\$?0[,.]6[–-]\$?2\b(?!\d)", "stale Home tier figure (0.6-2/M, superseded)"),
    (r"0[,.]77[–-]1[,.]20\b", "stale Cooperative tier figure (0.77-1.20/M, superseded)"),
    (r"0[,.]006[–-]\$?0[,.]24\b", "stale Home hourly figure (0.006-0.24/hr, superseded)"),
    (r"0[,.]046[–-]\$?0[,.]72\b", "stale Cooperative hourly figure (0.046-0.72/hr, superseded)"),
]

SELF_CORRECTION_MARKERS = [
    "Corrected 2026-08-13", "Resolved 2026-08-13", "Resolved (2026-08-13)",
    "Korjattu 2026-08-13", "Ratkaistu 2026-08-13", "Ratkaistu (2026-08-13)",
]

# Cells previously found to be broken, checked as fixed-value regression guards.
# (These are specific bug regression guards, not derivable from the CSV, because
# they check *storage type* and a *non-canonical intermediate assumption*
# [utilization], not a canonical output figure.)
XLSX_REGRESSION_CHECKS = {
    "hyperscale_utilization": {"sheet": "Hyperskaalataso", "cell": "B14", "expected": 0.60},
    "home_formula_label": {"sheet": "Kotitaloustaso", "cell": "C26"},
    "hyperscale_formula_label": {"sheet": "Hyperskaalataso", "cell": "C25"},
}


def check_dependencies():
    """Fail closed: report missing deps as failures, not silent skips."""
    problems = []
    available = {}
    for mod in REQUIRED_DEPS:
        try:
            importlib.import_module(mod)
            available[mod] = True
        except ImportError as e:
            available[mod] = False
            problems.append(
                f"DEPENDENCY MISSING: '{mod}' is not importable ({e}). "
                f"Run: pip install -r tools/requirements.txt in a working Python "
                f"environment. This checker treats a missing dependency as a "
                f"FAILURE, not a skipped check, so results stay trustworthy."
            )
    return problems, available


def load_canonical_figures():
    """Parse the canonical CSV and derive the exact figure-strings that must
    appear verbatim in the whitepaper, keyed by (tier, metric)."""
    rows = []
    with open(CANONICAL_CSV, newline="", encoding="utf-8") as f:
        for row in csv.DictReader(f):
            rows.append(row)
    if not rows:
        raise ValueError(f"{CANONICAL_CSV} parsed to zero rows — is the file empty or malformed?")

    # Canonical, full-layer, non-partial rows only: the check is "does the
    # canonical figure appear," not "do partial/sensitivity figures appear."
    required = []
    for row in rows:
        tier = row["tier"]
        if "partial" in tier.lower():
            continue  # e.g. Hyperscale_partial — deliberately not a canonical figure
        if "PARTIAL" in row.get("scope", ""):
            continue  # e.g. Professional tier — explicitly partial-layer, not required
        for bound in ("low", "high"):
            val = row.get(bound, "").strip()
            if val:
                # This repo's prose uses natural Finnish number formatting
                # (comma decimal separator) per TERMINOLOGY.md, while the
                # canonical CSV (shared, language-independent) stores plain
                # period-decimal values — convert before searching for the
                # figure verbatim in Finnish-language assets.
                required.append({
                    "tier": tier,
                    "metric": row["metric"],
                    "bound": bound,
                    "value": val.replace(".", ","),
                })
    return rows, required


def check_markdown_text(label, text, required_figures, is_whitepaper):
    """Run the stale-pattern and (for the whitepaper) canonical-figure checks
    against a block of text, returning a list of problem strings prefixed
    with the given label (a filename or 'PDF text')."""
    problems = []

    if is_whitepaper:
        for fig in required_figures:
            if fig["tier"] not in ("Home", "Cooperative"):
                continue  # Hyperscale/Professional/Retail_API figures vary in
                          # how they're formatted (ranges, mid-only, etc.) —
                          # Home/Cooperative low/high are the two blocker
                          # figures previously found stale, so those are the
                          # ones checked for verbatim presence package-wide.
            if fig["value"] not in text:
                problems.append(
                    f"{label}: canonical figure '{fig['value']}' ({fig['tier']} "
                    f"{fig['metric']} {fig['bound']}) not found anywhere — has "
                    f"the canonical range been removed, reworded, or not "
                    f"regenerated after a markdown edit?"
                )

    lines = text.split("\n")
    for i, line in enumerate(lines, start=1):
        for pattern, desc in STALE_PATTERNS:
            if re.search(pattern, line):
                if not any(marker in line for marker in SELF_CORRECTION_MARKERS):
                    problems.append(
                        f"{label}:{i}: found {desc} without a "
                        f"'Corrected'/'Resolved' self-correction marker on the same line"
                    )
    return problems


def check_markdown_files(required_figures):
    problems = []
    for relpath in MARKDOWN_FILES_TO_CHECK:
        path = os.path.join(ROOT, relpath)
        if not os.path.exists(path):
            problems.append(f"MISSING FILE expected for consistency check: {relpath}")
            continue
        with open(path, encoding="utf-8") as f:
            text = f.read()
        problems += check_markdown_text(relpath, text, required_figures, relpath == "01-whitepaper.md")
    return problems


def check_pdf(required_figures, deps_available):
    problems = []
    pdf_path = os.path.join(ROOT, "01-whitepaper.pdf")
    if not os.path.exists(pdf_path):
        return [f"MISSING FILE expected for consistency check: 01-whitepaper.pdf"]
    if not deps_available.get("pypdf"):
        return [
            "01-whitepaper.pdf: SKIPPED-AS-FAILURE — pypdf not available, cannot "
            "verify the built PDF matches the canonical model (see DEPENDENCY MISSING above)"
        ]

    from pypdf import PdfReader
    reader = PdfReader(pdf_path)
    text = "".join(page.extract_text() or "" for page in reader.pages)
    # PDF text extraction can insert stray whitespace inside numbers/words at
    # line-wrap points (a known artifact of this project's PDF pipeline) —
    # normalize runs of whitespace before matching so that doesn't produce
    # false positives.
    normalized = re.sub(r"\s+", " ", text)
    problems += check_markdown_text("01-whitepaper.pdf", normalized, required_figures, is_whitepaper=True)
    return problems


def check_xlsx(deps_available):
    problems = []
    xlsx_path = os.path.join(ROOT, "18-companion-data-model.xlsx")
    if not os.path.exists(xlsx_path):
        return [f"MISSING FILE expected for consistency check: 18-companion-data-model.xlsx"]
    if not deps_available.get("openpyxl"):
        return [
            "18-companion-data-model.xlsx: SKIPPED-AS-FAILURE — openpyxl not available, "
            "cannot verify workbook regression guards (see DEPENDENCY MISSING above)"
        ]

    import openpyxl
    wb = openpyxl.load_workbook(xlsx_path)

    check = XLSX_REGRESSION_CHECKS["hyperscale_utilization"]
    try:
        ws = wb[check["sheet"]]
        util = ws[check["cell"]].value
        if util is None or abs(float(util) - check["expected"]) > 1e-6:
            problems.append(
                f"18-companion-data-model.xlsx: {check['sheet']}!{check['cell']} utilization is "
                f"{util!r}, expected {check['expected']} (this paper's canonical mid-case) — "
                f"previously found set to 0.90, producing a wrong $0.091/M instead of $0.133/M"
            )
    except KeyError:
        problems.append(f"18-companion-data-model.xlsx: '{check['sheet']}' sheet not found")

    for key in ("home_formula_label", "hyperscale_formula_label"):
        check = XLSX_REGRESSION_CHECKS[key]
        try:
            ws = wb[check["sheet"]]
            cell = ws[check["cell"]]
            if cell.data_type == "f":
                problems.append(
                    f"18-companion-data-model.xlsx: {check['sheet']}!{check['cell']} is stored as a "
                    f"formula ({cell.value!r}) — this renders #NAME? in Excel; it should be a "
                    f"plain text cell"
                )
        except KeyError:
            problems.append(f"18-companion-data-model.xlsx: '{check['sheet']}' sheet not found")

    return problems


def check_pptx(required_figures, deps_available):
    problems = []
    pptx_path = os.path.join(ROOT, "19-slide-deck.pptx")
    if not os.path.exists(pptx_path):
        return [f"MISSING FILE expected for consistency check: 19-slide-deck.pptx"]
    if not deps_available.get("pptx"):
        return [
            "19-slide-deck.pptx: SKIPPED-AS-FAILURE — python-pptx not available, "
            "cannot verify deck text (see DEPENDENCY MISSING above)"
        ]

    from pptx import Presentation
    prs = Presentation(pptx_path)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        all_text.append(cell.text)
    full_text = "\n".join(all_text)

    for pattern, desc in STALE_PATTERNS:
        if re.search(pattern, full_text):
            problems.append(f"19-slide-deck.pptx: found {desc} in rendered slide text")

    for fig in required_figures:
        if fig["tier"] != "Home" or fig["metric"] != "cost_per_million_tokens":
            continue  # slide 11 cites the $/M-token figure specifically, not the derived hourly figures
        if fig["value"] not in full_text:
            problems.append(
                f"19-slide-deck.pptx: canonical Home-tier figure '{fig['value']}' not found in "
                f"rendered slide text — has slide 11 drifted from the canonical model?"
            )

    return problems


def run_check(name, fn, *args):
    """Run a check function, converting any exception into a reported
    problem instead of crashing the whole script — a broken runtime (e.g. a
    Python build with a non-functional pyexpat) must show up as a failed
    check, not an unhandled traceback that could be mistaken for 'the script
    didn't run, so there's nothing to worry about.'"""
    try:
        return fn(*args)
    except Exception as e:
        return [f"{name}: CHECK CRASHED ({type(e).__name__}: {e}) — treated as a failure, "
                f"not skipped. This usually means the Python environment is broken "
                f"(see tools/requirements.txt and the setup instructions in this "
                f"script's docstring) rather than a real content drift."]


def main():
    print("Checking canonical cost-model consistency...")
    print(f"  Canonical source: {os.path.relpath(CANONICAL_CSV, ROOT)}")

    all_problems = []

    dep_problems, deps_available = check_dependencies()
    all_problems += dep_problems

    try:
        _, required_figures = load_canonical_figures()
    except (OSError, ValueError) as e:
        print(f"\nFAIL — could not load canonical CSV: {e}")
        return 1

    all_problems += run_check("markdown files", check_markdown_files, required_figures)
    all_problems += run_check("01-whitepaper.pdf", check_pdf, required_figures, deps_available)
    all_problems += run_check("18-companion-data-model.xlsx", check_xlsx, deps_available)
    all_problems += run_check("19-slide-deck.pptx", check_pptx, required_figures, deps_available)

    print()
    if not all_problems:
        print("PASS — no drift detected across markdown, PDF, xlsx, and pptx assets. "
              "All required dependencies were available and all checks ran.")
        return 0

    print(f"FAIL — {len(all_problems)} consistency issue(s) found:\n")
    for p in all_problems:
        print(f"  - {p}")
    return 1


if __name__ == "__main__":
    sys.exit(main())
